"""
模块1：Document Loader
读取多类型课程资料，输出 List[Document]
支持：PDF / PPTX / DOCX / MD
"""

from dataclasses import dataclass
from pathlib import Path

from pypdf import PdfReader


SUPPORTED_SUFFIXES = {".pdf", ".pptx", ".docx", ".md"}


@dataclass
class Document:
    """单个文档单元（按页/按段/按幻灯片）"""

    content: str
    source: str  # 文件名
    page: int  # 页码/序号（从 1 开始）


def _validate_file(path: Path) -> None:
    if not path.exists():
        raise FileNotFoundError(f"文件不存在: {path}")
    if path.suffix.lower() not in SUPPORTED_SUFFIXES:
        raise ValueError(f"不支持的文件类型: {path.suffix}，支持: {sorted(SUPPORTED_SUFFIXES)}")


def load_pdf(file_path: str | Path) -> list[Document]:
    """加载单个 PDF，按页拆分为 Document 列表。"""
    path = Path(file_path)
    if not path.exists():
        raise FileNotFoundError(f"PDF 不存在: {path}")
    if path.suffix.lower() != ".pdf":
        raise ValueError(f"非 PDF 文件: {path}")

    reader = PdfReader(path)
    documents: list[Document] = []
    for i, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        documents.append(Document(content=text.strip(), source=path.name, page=i))
    return documents


def load_markdown(file_path: str | Path) -> list[Document]:
    """加载单个 Markdown 文件，作为一个 Document 返回。"""
    path = Path(file_path)
    if not path.exists():
        raise FileNotFoundError(f"Markdown 不存在: {path}")
    if path.suffix.lower() != ".md":
        raise ValueError(f"非 Markdown 文件: {path}")
    text = path.read_text(encoding="utf-8")
    return [Document(content=text.strip(), source=path.name, page=1)]


def load_docx(file_path: str | Path) -> list[Document]:
    """加载单个 Word（.docx）文件，作为一个 Document 返回。"""
    path = Path(file_path)
    if not path.exists():
        raise FileNotFoundError(f"Word 文件不存在: {path}")
    if path.suffix.lower() != ".docx":
        raise ValueError(f"非 DOCX 文件: {path}")
    try:
        from docx import Document as DocxDocument  # type: ignore[import-not-found]
    except Exception as exc:  # pragma: no cover - 依赖环境相关
        raise RuntimeError("请先安装 python-docx 以支持 .docx 加载") from exc

    doc = DocxDocument(path)
    lines = [p.text.strip() for p in doc.paragraphs if p.text and p.text.strip()]
    return [Document(content="\n".join(lines), source=path.name, page=1)]


def load_pptx(file_path: str | Path) -> list[Document]:
    """加载单个 PowerPoint（.pptx）文件，按幻灯片拆分为 Document 列表。"""
    path = Path(file_path)
    if not path.exists():
        raise FileNotFoundError(f"PPT 文件不存在: {path}")
    if path.suffix.lower() != ".pptx":
        raise ValueError(f"非 PPTX 文件: {path}")
    try:
        from pptx import Presentation  # type: ignore[import-not-found]
    except Exception as exc:  # pragma: no cover - 依赖环境相关
        raise RuntimeError("请先安装 python-pptx 以支持 .pptx 加载") from exc

    prs = Presentation(str(path))
    docs: list[Document] = []
    for slide_idx, slide in enumerate(prs.slides, start=1):
        texts: list[str] = []
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if text and str(text).strip():
                texts.append(str(text).strip())
        docs.append(Document(content="\n".join(texts), source=path.name, page=slide_idx))
    return docs


def load_document(file_path: str | Path) -> list[Document]:
    """统一单文件入口，按扩展名分发到具体加载函数。"""
    path = Path(file_path)
    _validate_file(path)
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return load_pdf(path)
    if suffix == ".md":
        return load_markdown(path)
    if suffix == ".docx":
        return load_docx(path)
    if suffix == ".pptx":
        return load_pptx(path)
    raise ValueError(f"不支持的文件类型: {suffix}")


def load_documents_from_dir(
    dir_path: str | Path,
    *,
    allowed_suffixes: set[str] | None = None,
) -> list[Document]:
    """
    加载目录下所有受支持文档（默认支持 pdf/pptx/docx/md）。
    """
    p = Path(dir_path)
    if not p.is_dir():
        raise NotADirectoryError(f"目录不存在: {p}")

    suffixes = {s.lower() for s in (allowed_suffixes or SUPPORTED_SUFFIXES)}
    file_paths = sorted([fp for fp in p.iterdir() if fp.is_file() and fp.suffix.lower() in suffixes])
    all_docs: list[Document] = []
    for fp in file_paths:
        all_docs.extend(load_document(fp))
    return all_docs


def load_pdfs_from_dir(dir_path: str | Path) -> list[Document]:
    """
    兼容旧接口：仅加载目录下 PDF 文件。
    """
    return load_documents_from_dir(dir_path, allowed_suffixes={".pdf"})


if __name__ == "__main__":
    docs = load_documents_from_dir("data")
    print(f"共加载 {len(docs)} 条文档单元")
    for d in docs[:5]:
        print(f"  [{d.source} p.{d.page}] {d.content[:80]}...")
